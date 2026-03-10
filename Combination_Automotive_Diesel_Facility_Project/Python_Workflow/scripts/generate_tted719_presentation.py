from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

out = Path(
    "Combination_Automotive_Diesel_Facility_Project/Python_Workflow/outputs/Training_Facility_Drawings_v1.1/TTED719_Presentation_Rodriguez_v1.1.pptx"
)
out.parent.mkdir(parents=True, exist_ok=True)


def add_slide(title, bullets, notes):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = notes


# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Training Facility Design — Rodriguez (TTED 719)"
slide.shapes.placeholders[1].text = (
    "10–15 minute overview: layout, systems, cost, and reproducibility"
)
slide.notes_slide.notes_text_frame.text = "Speaker: Brief self-intro, state scope: CAD sources, inspection-ready PDFs, equipment & cost, maintenance, and reproducible scripts. "

# Slides content
add_slide(
    "Presentation Overview",
    [
        "Scope & objectives",
        "Key deliverables included",
        "What reviewers should look for",
    ],
    "Speaker: State objectives: present plan, show evidence of compliance to TTED 719, point reviewers to manifest.md and checksums.sha256 for authoritative files.",
)

add_slide(
    "Authoritative CAD Sources",
    [
        "Master layered DXF: training_facility_plan_layered.dxf",
        "Derived discipline DXFs and labeled variants",
        "Print-ready: facility_layout_full.pdf",
    ],
    "Speaker: Explain layered DXF role as single source of truth; changes propagate to PDFs via Python_Workflow scripts. Mention facility_layout_labeled.dxf and engineering variant for discipline use.",
)

add_slide(
    "Circulation & Bay Layout",
    [
        "Student bays, heavy-duty and light-duty zones",
        "EV charging & service bays allocation",
        "Circulation routes for safety and instruction",
    ],
    "Speaker: Walk through bay counts and placements; point to labeled plan and legend for bay IDs. Emphasize instructional flow and emergency egress considerations.",
)

add_slide(
    "Equipment & Furniture",
    [
        "Essential equipment by bay (CSV/XLSX)",
        "Furniture sizing and placement (FURN_Plan.dxf)",
        "Procurement workbook attached",
    ],
    "Speaker: Highlight equipment_bay_mapping.xlsx and Vendor_Procurement_Master.xlsx. Explain how costs roll up and link to totalcost_per_bay.png.",
)

add_slide(
    "Safety, HVAC & Utilities",
    [
        "Exhaust, ventilation placeholders noted",
        "Lighting & footcandle considerations",
        "Electrical distribution & EV infrastructure",
    ],
    "Speaker: Call out engineering annotations and legend. Note mechanical_services.csv and electrical_loads.csv for system assumptions.",
)

add_slide(
    "Cost & Maintenance Overview",
    [
        "Total facility estimate summary",
        "Maintenance and replacement plan included",
        "Consumables and lifecycle notes",
    ],
    "Speaker: Refer to Vendor_Procurement_Master and maintenance tabs. Emphasize transparent assumptions and reproducible calculations.",
)

add_slide(
    "Inspection-Ready Deliverables",
    [
        "Print-ready PDFs and portfolio package",
        "Manifest + checksums for integrity",
        "Submission ZIP available",
    ],
    "Speaker: Tell reviewers how to verify: open facility_layout_full.pdf, compare hashes in checksums.sha256, see manifest.md for file roles.",
)

add_slide(
    "Reproducibility & Automation",
    [
        "Python_Workflow scripts regenerate artifacts",
        "CI checks (mypy/black) and typed code",
        "How to run: run_pipeline.ps1 / scripts entrypoints",
    ],
    "Speaker: Explain that scripts automate DXF→PDF, legend assembly, cost PDF generation; mention .github workflow for type checks.",
)

add_slide(
    "3D Visualization & Appendix",
    [
        "Canonical Blender model included",
        "Rendered views included in appendix PDFs",
        "Optional STEP/STL export available",
    ],
    "Speaker: Show a couple of rendered images (in portfolio). Explain 3D model helps stakeholders visualize layout and adjacency.",
)

add_slide(
    "Next Steps & Requests",
    [
        "Questions for reviewers",
        "Requested sign-offs or clarifications",
        "How to regenerate deliverables locally",
    ],
    "Speaker: Invite reviewers to request additional discipline-level PDFs (ELEC/MECH), or accept changes. Provide brief run instructions: activate venv, pip install -r requirements, run scripts.",
)

add_slide(
    "Q&A",
    [
        "Contact and repo pointers",
        "manifest.md and TTED719_mapping.pdf for quick review",
    ],
    "Speaker: Close with contact info and pointer to TTED719_mapping.pdf and Instructor Review Checklist. Thank the reviewers.",
)

prs.save(str(out))
print(f"Saved presentation: {out}")
